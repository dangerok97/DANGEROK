"""Generate synthetic (non-personal) fixtures for intelligent documents verification.

Run:  backend/.venv/Scripts/python.exe tests/fixtures/intel_docs/generate_fixtures.py
"""
from __future__ import annotations

import csv
import io
from pathlib import Path

ROOT = Path(__file__).resolve().parent


VISITA = """\
CONFERMA VISITA SPECIALISTICA
Struttura: Poliambulatorio Nord Demo
Indirizzo: Via dei Campioni 12
Città: Milano
Data: 18 settembre 2027
Ora: 10:30
Istruzioni: Presentarsi 15 minuti prima con tessera sanitaria.
Codice prenotazione: PREN-DEMO-4411
Nota: documento sintetico di test — nessun dato reale.
"""

CONCERTO = """\
BIGLIETTO CONCERTO
Artista: Aurora Nova
Data: 22 ottobre 2027
Apertura porte: 19:30
Inizio evento: 21:00
Venue: Arena Demo
Indirizzo: Piazza delle Prove 3
Città: Bologna
Codice ordine: ORD-DEMO-7788
"""

TRENO = """\
BIGLIETTO FERROVIARIO
Stazione partenza: Firenze Santa Maria Novella
Stazione arrivo: Roma Termini
Data: 05 novembre 2027
Orario partenza: 08:15
Orario arrivo: 09:52
Codice prenotazione: PNR-DEMO-3322
"""

DISPENSA = """\
Dispensa universitaria sintetica
Materia: Antropologia culturale
Argomento: Habitus e campo in Bourdieu
Capitolo 1
Definizione: L'habitus è un sistema di disposizioni durature e trasmissibili.
Concetto: capitale culturale
Concetto: campo sociale
Concetto: violenza simbolica
Lezione: le pratiche quotidiane riproducono strutture sociali.
Bibliografia: Bourdieu, Outline of a Theory of Practice.
Domanda di ripasso: Cos'è l'habitus?
Domanda di ripasso: Come si collega al campo?
"""

ADMIN = """\
Comunicazione amministrativa
Mittente: Ufficio Protocollo Demo
Oggetto: Richiesta integrazione documentazione
Azione richiesta: Caricare il modulo aggiornato sul portale
Scadenza: 30 aprile 2027
Riferimenti: Prot. DEMO/2027/9988
"""

AMBIGUA = """\
Appuntamento generico
Data: 03/04/2027
Ora: 10:30
Luogo: Ufficio
Contesto insufficiente per disambiguare giorno/mese.
"""


def _write_text(name: str, body: str) -> Path:
    p = ROOT / name
    p.write_text(body.strip() + "\n", encoding="utf-8")
    return p


def _png(text: str, name: str, *, rotate: int = 0, blurry: bool = False) -> Path:
    from PIL import Image, ImageDraw, ImageFilter, ImageFont

    img = Image.new("RGB", (1100, 320), "white")
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 28)
    except Exception:
        font = ImageFont.load_default()
    y = 40
    for line in text.splitlines():
        d.text((30, y), line[:90], fill="black", font=font)
        y += 36
    if rotate:
        img = img.rotate(rotate, expand=True, fillcolor="white")
    if blurry:
        img = img.filter(ImageFilter.GaussianBlur(radius=4))
        img = img.resize((img.width // 3, img.height // 3)).resize(img.size)
    path = ROOT / name
    img.save(path, format="PNG" if name.lower().endswith(".png") else "JPEG", quality=40 if blurry else 90)
    return path


def _scanned_pdf(name: str, text: str) -> Path:
    from PIL import Image, ImageDraw, ImageFont
    import pypdfium2 as pdfium  # noqa: F401 — ensure available
    # Build image pages then wrap with pypdf if needed — simpler: reportlab-less via img→pdf with pillow + img2pdf or pypdf
    from pypdf import PdfWriter
    import io as _io

    pages = []
    for chunk in (text[:300], text[300:600] or "Pagina 2 sintetica"):
        img = Image.new("RGB", (900, 1200), "white")
        d = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 22)
        except Exception:
            font = ImageFont.load_default()
        y = 60
        for line in chunk.splitlines() or ["pagina"]:
            d.text((40, y), line[:70], fill="black", font=font)
            y += 30
        buf = _io.BytesIO()
        img.save(buf, format="PDF")
        pages.append(buf.getvalue())

    writer = PdfWriter()
    for blob in pages:
        writer.append(_io.BytesIO(blob))
    out = ROOT / name
    with out.open("wb") as f:
        writer.write(f)
    return out


def _docx(name: str, body: str) -> Path:
    from docx import Document

    doc = Document()
    for line in body.strip().splitlines():
        doc.add_paragraph(line)
    path = ROOT / name
    doc.save(path)
    return path


def _pptx(name: str, body: str) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dispensa sintetica"
    slide.placeholders[1].text = body[:500]
    path = ROOT / name
    prs.save(path)
    return path


def _csv(name: str) -> Path:
    path = ROOT / name
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["campo", "valore"])
        w.writerow(["evento", "Workshop Demo"])
        w.writerow(["data", "10 gennaio 2028"])
        w.writerow(["citta", "Torino"])
    return path


def _md(name: str, body: str) -> Path:
    return _write_text(name, "# Documento markdown sintetico\n\n" + body)


def _text_pdf(name: str, body: str) -> Path:
    # Minimal text PDF via pypdf + page content is hard; use reportlab-free approach:
    # write as plain and also a simple PDF with page extractable text using pypdf canvas alternative.
    try:
        from pypdf import PdfWriter
        from pypdf.generic import DictionaryObject, DecodedStreamObject, NameObject, NumberObject, ArrayObject

        # Fallback: store as .txt companion and create image-based if needed.
        # Prefer writing with a tiny handwritten PDF stream.
    except Exception:
        pass
    # Use pillow PDF (image) for reliability in tests that OCR; for text PDF use fpdf if available.
    try:
        from fpdf import FPDF  # type: ignore

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        for line in body.splitlines():
            pdf.multi_cell(0, 8, line[:100])
        path = ROOT / name
        pdf.output(str(path))
        return path
    except Exception:
        return _scanned_pdf(name, body)


def main() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)
    outs = [
        _write_text("caso_a_visita.txt", VISITA),
        _write_text("caso_b_concerto.txt", CONCERTO),
        _write_text("caso_c_treno.txt", TRENO),
        _write_text("caso_d_dispensa.txt", DISPENSA),
        _write_text("caso_e_admin.txt", ADMIN),
        _write_text("caso_f_ambigua.txt", AMBIGUA),
        _png(VISITA, "ocr_readable.png"),
        _png(VISITA, "ocr_tilted.jpg", rotate=12),
        _png("DATA 12/03/2027 ORA 15:30 INDIRIZZO Via Demo 1 Milano", "ocr_datetime_address.png"),
        _png("xx##@@ quasi illeggibile", "ocr_low_quality.png", blurry=True),
        _scanned_pdf("ocr_scanned.pdf", VISITA),
        _docx("formato_visita.docx", VISITA),
        _pptx("formato_dispensa.pptx", DISPENSA),
        _csv("formato_eventi.csv"),
        _md("formato_notes.md", DISPENSA),
        _write_text("formato_plain.txt", CONCERTO),
        _text_pdf("formato_testuale.pdf", CONCERTO),
    ]
    print(f"generated {len(outs)} fixtures in {ROOT}")


if __name__ == "__main__":
    main()
