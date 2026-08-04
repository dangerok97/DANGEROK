"""Local text extraction for DOCX / PPTX / CSV / Markdown."""
from __future__ import annotations

import csv
import io
import logging
from typing import Optional

from documents.extraction import ExtractionResult, clean_text, detect_language

logger = logging.getLogger("ora.documents.office")


def extract_docx(blob: bytes) -> ExtractionResult:
    r = ExtractionResult(engine="python-docx")
    try:
        from docx import Document
    except Exception:
        r.error_code = "docx_lib_missing"
        return r
    try:
        doc = Document(io.BytesIO(blob))
        parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        r.text = clean_text("\n".join(parts))
        r.pages = max(1, len(parts) // 40)
        r.language = detect_language(r.text)
    except Exception as e:
        r.error_code = "docx_extract_failed"
        r.warnings.append(type(e).__name__)
    return r


def extract_pptx(blob: bytes) -> ExtractionResult:
    r = ExtractionResult(engine="python-pptx")
    try:
        from pptx import Presentation
    except Exception:
        r.error_code = "pptx_lib_missing"
        return r
    try:
        prs = Presentation(io.BytesIO(blob))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_bits: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        slide_bits.append(t)
            if slide_bits:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_bits))
        r.text = clean_text("\n\n".join(parts))
        r.pages = len(prs.slides)
        r.language = detect_language(r.text)
    except Exception as e:
        r.error_code = "pptx_extract_failed"
        r.warnings.append(type(e).__name__)
    return r


def extract_by_mime(blob: bytes, mime_type: str) -> Optional[ExtractionResult]:
    mime = (mime_type or "").lower().split(";")[0].strip()
    if mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return extract_docx(blob)
    if mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return extract_pptx(blob)
    return None
